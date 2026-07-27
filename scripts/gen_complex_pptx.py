#!/usr/bin/env python3
"""Generate a complex PPTX for LibreOffice / filebox Office preview E2E tests."""

from __future__ import annotations

import argparse
import io
import random
import string
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


def make_placeholder_image(width: int, height: int, label: str) -> io.BytesIO:
    img = Image.new("RGB", (width, height), color=(30 + random.randint(0, 40), 80, 140))
    draw = ImageDraw.Draw(img)
    draw.rectangle((10, 10, width - 10, height - 10), outline=(255, 255, 255), width=3)
    draw.text((24, height // 2 - 12), label, fill=(255, 255, 255))
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    buf.seek(0)
    return buf


def lorem(n: int = 40) -> str:
    words = [
        "analysis", "benchmark", "cluster", "dataset", "endpoint", "forecast",
        "gradient", "histogram", "iteration", "kernel", "latency", "matrix",
        "network", "outlier", "pipeline", "quantile", "regression", "sample",
        "throughput", "variance", "workflow", "yield", "zone",
    ]
    return " ".join(random.choice(words) for _ in range(n))


def add_title_slide(prs: Presentation, title: str, subtitle: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = title
    slide.placeholders[1].text = subtitle


def add_bullet_slide(prs: Presentation, title: str, bullets: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, bullet in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = bullet
        p.level = i % 3


def add_image_slide(prs: Presentation, title: str, label: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    pic = make_placeholder_image(960, 540, label)
    slide.shapes.add_picture(pic, Inches(1.0), Inches(1.5), width=Inches(8.0))


def add_table_slide(prs: Presentation, title: str, rows: int, cols: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    table = slide.shapes.add_table(rows, cols, Inches(0.8), Inches(1.6), Inches(8.4), Inches(4.5)).table
    for r in range(rows):
        for c in range(cols):
            cell = table.cell(r, c)
            if r == 0:
                cell.text = f"Col {c + 1}"
            else:
                cell.text = f"{random.randint(1, 9999):04d}"


def add_chart_slide(prs: Presentation, title: str, series_count: int) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    chart_data = CategoryChartData()
    chart_data.categories = [f"Q{i}" for i in range(1, 9)]
    for s in range(series_count):
        chart_data.add_series(f"Series {s + 1}", [random.randint(10, 100) for _ in range(8)])
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1.0),
        Inches(1.5),
        Inches(8.0),
        Inches(4.8),
        chart_data,
    )


def add_textbox_slide(prs: Presentation, title: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    slide.shapes.title.text = title
    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.4), Inches(8.4), Inches(5.0))
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = lorem(80)
    p.font.size = Pt(14)
    p.alignment = PP_ALIGN.JUSTIFY
    p2 = tf.add_paragraph()
    p2.text = lorem(60)
    p2.font.size = Pt(12)
    p2.alignment = PP_ALIGN.LEFT


def build_presentation(slide_count: int) -> Presentation:
    prs = Presentation()
    add_title_slide(
        prs,
        "filebox Complex PPTX E2E",
        f"{slide_count} slides — images, tables, charts, dense text",
    )

    builders = [
        lambda i: add_bullet_slide(prs, f"Agenda {i}", [f"Point {j}: {lorem(8)}" for j in range(1, 6)]),
        lambda i: add_image_slide(prs, f"Diagram {i}", f"IMG-{i:03d}"),
        lambda i: add_table_slide(prs, f"Metrics {i}", 6, 5),
        lambda i: add_chart_slide(prs, f"Trend {i}", 3),
        lambda i: add_textbox_slide(prs, f"Notes {i}"),
    ]

    # Fill remaining slides cycling through layouts.
    for i in range(1, slide_count):
        builders[i % len(builders)](i)

    return prs


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("-o", "--output", type=Path, required=True)
    parser.add_argument("-n", "--slides", type=int, default=70)
    args = parser.parse_args()

    random.seed(42)
    prs = build_presentation(args.slides)
    args.output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(args.output))
    size = args.output.stat().st_size
    print(f"Wrote {args.output} ({args.slides} slides, {size:,} bytes)")


if __name__ == "__main__":
    main()
